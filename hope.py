"""HOPE Chatbot (multi-doc, PPT/PDF/Excel-aware)

What this does:
- Loads text from one or more documents (PPTX, PDF, TXT/MD, CSV, XLSX)
- Chunks it
- Builds embeddings for each chunk
- On each user question, retrieves the most relevant chunks
- Sends the chunks + question to the model to answer grounded in the documents

Setup:
  pip install openai python-pptx numpy pypdf openpyxl

Examples:
  export OPENAI_API_KEY="..."

  # Single deck
  python hope.py --files "/path/to/presentation.pptx"

  # Multiple sources
  python hope.py --files "/path/to/slides.pptx" "/path/to/handout.pdf" "/path/to/data.xlsx" "/path/to/notes.txt"

Notes:
- This is a simple in-memory RAG. For big corpora, use a vector DB.
- Answers are constrained to retrieved context; if context doesn't contain it, HOPE will say so.
"""

import os
import argparse
from typing import List
import csv
from pathlib import Path
from dotenv import load_dotenv
load_dotenv()
from pypdf import PdfReader
from openpyxl import load_workbook

# ... (other imports as needed)

# Placeholder for DocChunk and InMemoryVectorStore classes and embed_texts, get_client, chunk_text, etc.
# You would place the rest of your chatbot code and class definitions here.

def extract_text_from_pptx_file(pptx_path: str) -> str:
    """Extract text from a PowerPoint file using python-pptx."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = (shape.text or "").strip()
                if txt:
                    slide_text.append(txt)
        if slide_text:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
    return "\n\n".join(parts).strip()


def extract_text_from_pdf_file(pdf_path: str) -> str:
    """Extract text from a PDF using pypdf."""
    reader = PdfReader(pdf_path)
    parts: List[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = (page.extract_text() or "").strip()
        except Exception:
            txt = ""
        if txt:
            parts.append(f"[Page {i}]\n{txt}")
    return "\n\n".join(parts).strip()


def extract_text_from_txt_file(txt_path: str) -> str:
    """Extract text from a plain text / markdown file."""
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extract_text_from_csv_file(csv_path: str, max_rows: int = 5000) -> str:
    """Extract text from a CSV by rendering rows as pipe-delimited lines."""
    parts: List[str] = []
    with open(csv_path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader, start=1):
            if i > max_rows:
                parts.append(f"[Truncated after {max_rows} rows]")
                break
            line = " | ".join(cell.strip() for cell in row if cell is not None)
            if line.strip():
                parts.append(line)
    return "\n".join(parts).strip()


def extract_text_from_xlsx_file(xlsx_path: str, max_cells: int = 20000) -> str:
    """Extract text from an Excel workbook (values only)."""
    wb = load_workbook(xlsx_path, data_only=True, read_only=True)
    parts: List[str] = []
    cell_count = 0

    for ws in wb.worksheets:
        parts.append(f"[Sheet] {ws.title}")
        for row in ws.iter_rows(values_only=True):
            if cell_count >= max_cells:
                parts.append(f"[Truncated after {max_cells} cells]")
                wb.close()
                return "\n".join(parts).strip()

            # Render a row only if it contains something non-empty
            rendered: List[str] = []
            for v in row:
                cell_count += 1
                if v is None:
                    rendered.append("")
                else:
                    rendered.append(str(v).strip())

            if any(x for x in rendered):
                parts.append(" | ".join(rendered))

        parts.append("")  # blank line between sheets

    wb.close()
    return "\n".join(parts).strip()


# ... (chunk_text and other functions/classes)

def load_documents(file_paths: List[str]) -> str:
    """Load and concatenate text from supported document formats."""
    supported = {".pptx", ".pdf", ".txt", ".md", ".csv", ".xlsx"}
    outputs: List[str] = []

    for p in file_paths:
        path = Path(p)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {p}")

        ext = path.suffix.lower()
        if ext not in supported:
            raise ValueError(
                f"Unsupported file type: {ext} for {p}. Supported: {sorted(supported)}"
            )

        header = f"[Source {path.name}]"
        if ext == ".pptx":
            body = extract_text_from_pptx_file(str(path))
        elif ext == ".pdf":
            body = extract_text_from_pdf_file(str(path))
        elif ext in {".txt", ".md"}:
            body = extract_text_from_txt_file(str(path))
        elif ext == ".csv":
            body = extract_text_from_csv_file(str(path))
        elif ext == ".xlsx":
            body = extract_text_from_xlsx_file(str(path))
        else:
            body = ""

        if body:
            outputs.append(f"{header}\n{body}")

    return "\n\n=====\n\n".join(outputs).strip()


# ... (embed_texts, chunk_text, DocChunk, InMemoryVectorStore, etc.)

def build_store_from_files(client, file_paths: List[str]):
    text = load_documents(file_paths)
    if not text:
        return InMemoryVectorStore([])

    chunks = chunk_text(text)
    embs = embed_texts(client, chunks)

    doc_chunks: List[DocChunk] = [DocChunk(t, e) for t, e in zip(chunks, embs)]
    return InMemoryVectorStore(doc_chunks)


def main():
    parser = argparse.ArgumentParser(description="HOPE Chatbot (multi-doc)")
    parser.add_argument(
        "--files",
        nargs="+",
        required=True,
        help="One or more input files (.pptx, .pdf, .txt, .md, .csv, .xlsx)",
    )
    parser.add_argument("--top_k", type=int, default=5, help="How many chunks to retrieve")
    args = parser.parse_args()

    for fpath in args.files:
        if not os.path.exists(fpath):
            print(f"Error: File not found: {fpath}")
            return 1

    client = get_client()

    print("Loading and indexing documents:")
    for f in args.files:
        print(f"- {f}")

    store = build_store_from_files(client, args.files)

    # ... (rest of your chatbot main loop)


if __name__ == "__main__":
    main()
 